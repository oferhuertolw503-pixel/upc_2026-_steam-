"""
Steam游戏市场数据分析 — 演讲PPT生成器
======================================
依赖: pip install python-pptx

运行: python generate_ppt.py
输出: Steam游戏市场数据分析_演讲.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ============================================================
# 全局设置
# ============================================================
prs = Presentation()
prs.slide_width = Inches(13.333)   # 16:9 宽屏
prs.slide_height = Inches(7.5)

# 配色方案
DARK_BLUE = RGBColor(0x0D, 0x2B, 0x4E)
MEDIUM_BLUE = RGBColor(0x1A, 0x56, 0x9A)
LIGHT_BLUE = RGBColor(0x2E, 0x86, 0xC1)
ACCENT_ORANGE = RGBColor(0xE6, 0x7E, 0x22)
ACCENT_GREEN = RGBColor(0x27, 0xAE, 0x60)
ACCENT_RED = RGBColor(0xC0, 0x39, 0x2B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xEC, 0xF0, 0xF1)
DARK_GRAY = RGBColor(0x2C, 0x3E, 0x50)
MEDIUM_GRAY = RGBColor(0x7F, 0x8C, 0x8D)
BOX_BG = RGBColor(0xF7, 0xF9, 0xFC)

# ============================================================
# 辅助函数
# ============================================================

def add_bg(slide, color=DARK_BLUE):
    """设置幻灯片纯色背景"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, color, alpha=None):
    """添加矩形色块"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT,
                font_name='Arial', line_spacing=1.2):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(line_spacing * font_size - font_size)
    return txBox

def add_multiline_textbox(slide, left, top, width, height, lines, default_size=16):
    """添加多行文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if isinstance(line_data, str):
            p.text = line_data
            p.font.size = Pt(default_size)
            p.font.color.rgb = DARK_GRAY
        elif isinstance(line_data, dict):
            p.text = line_data.get('text', '')
            p.font.size = Pt(line_data.get('size', default_size))
            p.font.color.rgb = line_data.get('color', DARK_GRAY)
            p.font.bold = line_data.get('bold', False)
            p.font.name = line_data.get('font', 'Arial')
            if 'alignment' in line_data:
                p.alignment = line_data['alignment']
        p.space_after = Pt(4)
    return txBox

def add_image_safe(slide, image_path, left, top, width, height=None):
    """安全添加图片"""
    if os.path.exists(image_path):
        if height:
            return slide.shapes.add_picture(image_path, left, top, width, height)
        else:
            return slide.shapes.add_picture(image_path, left, top, width=width)
    return None

def slide_title_bar(slide, title, subtitle=None):
    """统一的标题栏样式"""
    # 顶部色条
    add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), ACCENT_ORANGE)
    # 深蓝顶栏
    add_rect(slide, Inches(0), Inches(0.08), prs.slide_width, Inches(1.1), DARK_BLUE)
    # 标题文字
    add_textbox(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.7),
                title, font_size=32, color=WHITE, bold=True)
    if subtitle:
        add_textbox(slide, Inches(0.8), Inches(0.78), Inches(11), Inches(0.4),
                    subtitle, font_size=14, color=RGBColor(0xBD, 0xC3, 0xC7))
    # 底部页码条
    add_rect(slide, Inches(0), Inches(7.35), prs.slide_width, Inches(0.15), DARK_BLUE)

def add_page_number(slide, num):
    add_textbox(slide, Inches(12.4), Inches(7.15), Inches(0.8), Inches(0.3),
                str(num), font_size=10, color=MEDIUM_GRAY, alignment=PP_ALIGN.RIGHT)

def add_kpi_box(slide, left, top, width, height, number, label, color=MEDIUM_BLUE):
    """添加 KPI 指标卡片"""
    shape = add_rect(slide, left, top, width, height, BOX_BG)
    shape.shadow.inherit = False

    # 顶部色条
    add_rect(slide, left, top, width, Inches(0.06), color)

    add_textbox(slide, left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.5),
                number, font_size=28, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.1), top + Inches(0.7), width - Inches(0.2), Inches(0.4),
                label, font_size=11, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

def add_icon_bullet(slide, left, top, width, height, icon_text, title, desc, color=MEDIUM_BLUE):
    """添加带图标的要点"""
    # 圆形图标
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top + Inches(0.05), Inches(0.45), Inches(0.45))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = icon_text
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    add_textbox(slide, left + Inches(0.6), top, width - Inches(0.6), Inches(0.35),
                title, font_size=15, color=DARK_GRAY, bold=True)
    add_textbox(slide, left + Inches(0.6), top + Inches(0.35), width - Inches(0.6), Inches(0.55),
                desc, font_size=11, color=MEDIUM_GRAY)

# ============================================================
# 第1页：封面
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK_BLUE)

# 装饰线
add_rect(slide, Inches(1.2), Inches(2.6), Inches(0.8), Inches(0.06), ACCENT_ORANGE)

add_textbox(slide, Inches(1.2), Inches(2.8), Inches(11), Inches(1.2),
            'Steam 游戏市场数据分析',
            font_size=44, color=WHITE, bold=True)

add_textbox(slide, Inches(1.2), Inches(4.0), Inches(10), Inches(0.6),
            '基于 Steam Video Games 2024 数据集的系统分析与洞察',
            font_size=20, color=RGBColor(0xBD, 0xC3, 0xC7))

add_textbox(slide, Inches(1.2), Inches(5.6), Inches(5), Inches(0.4),
            '数据思维与人工智能 · 课程大作业',
            font_size=14, color=MEDIUM_GRAY)

add_textbox(slide, Inches(1.2), Inches(6.1), Inches(5), Inches(0.4),
            '2026年6月',
            font_size=14, color=MEDIUM_GRAY)

# ============================================================
# 第2页：目录
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_title_bar(slide, '汇报大纲', 'AGENDA')
add_page_number(slide, 2)

agenda = [
    ('01', '项目背景与数据集', '分析动机、数据来源与规模'),
    ('02', '分析框架与方法论', '四维度框架 + 方法选择优势'),
    ('03', '维度一：综合评价', '价格/评分/销量分布、相关性矩阵'),
    ('04', '维度二：定价策略', '价格区间对比、免费vs付费'),
    ('05', '维度三：类型趋势', '市场份额、评分对比、年度趋势'),
    ('06', '维度四：玩家评价', '影响因素探索、Metacritic对比'),
    ('07', '结论与建议', '核心发现、对开发者的启示'),
]

for i, (num, title, desc) in enumerate(agenda):
    y = Inches(1.6) + Inches(0.75) * i
    add_icon_bullet(slide, Inches(1.5), y, Inches(10), Inches(0.6),
                    num, title, desc, color=[DARK_BLUE, MEDIUM_BLUE, LIGHT_BLUE,
                                             MEDIUM_BLUE, LIGHT_BLUE, MEDIUM_BLUE,
                                             ACCENT_ORANGE][i])

# ============================================================
# 第3页：项目背景与数据集
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_title_bar(slide, '项目背景与数据集概述', 'BACKGROUND & DATASET')
add_page_number(slide, 3)

# 左侧-背景
add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.4),
            '为什么分析 Steam？', font_size=20, color=DARK_BLUE, bold=True)
add_multiline_textbox(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(2.5), [
    {'text': 'Steam 是全球最大 PC 游戏平台', 'size': 14, 'bold': True, 'color': DARK_GRAY},
    {'text': '', 'size': 6},
    {'text': '• 月活跃用户超过 1.3 亿', 'size': 13, 'color': DARK_GRAY},
    {'text': '• 平台拥有 97,000+ 款游戏和应用', 'size': 13, 'color': DARK_GRAY},
    {'text': '• 支持 28 种语言、全球 60+ 支付方式', 'size': 13, 'color': DARK_GRAY},
    {'text': '', 'size': 6},
    {'text': '• 海量公开数据为学术分析提供素材', 'size': 13, 'color': DARK_GRAY},
    {'text': '• 游戏产业数据驱动决策的典型场景', 'size': 13, 'color': DARK_GRAY},
    {'text': '• 分析结果对开发者有直接商业价值', 'size': 13, 'color': DARK_GRAY},
])

# 右侧-KPI卡片
add_kpi_box(slide, Inches(7.3), Inches(1.5), Inches(2.6), Inches(1.1),
            '97,428', '游戏总数', DARK_BLUE)
add_kpi_box(slide, Inches(10.1), Inches(1.5), Inches(2.6), Inches(1.1),
            '40', '数据字段', MEDIUM_BLUE)
add_kpi_box(slide, Inches(7.3), Inches(2.8), Inches(2.6), Inches(1.1),
            '79.8%', '付费游戏占比', ACCENT_GREEN)
add_kpi_box(slide, Inches(10.1), Inches(2.8), Inches(2.6), Inches(1.1),
            '2006-2025', '时间跨度', LIGHT_BLUE)

# 底部-数据集说明
add_textbox(slide, Inches(0.8), Inches(4.5), Inches(11.5), Inches(0.4),
            '数据来源：Kaggle - Steam Video Games 2024', font_size=13, color=DARK_BLUE, bold=True)
add_multiline_textbox(slide, Inches(0.8), Inches(5.0), Inches(11.5), Inches(2), [
    {'text': '关键字段：价格、估计拥有者、好评/差评数、Metacritic评分、类型标签、发行日期、成就数、DLC数、支持平台等', 'size': 12, 'color': DARK_GRAY},
    {'text': '数据清洗：解析拥有者区间为中位数值、提取发行年份、计算好评率、展开多类型标签、构建价格/拥有者分类变量', 'size': 12, 'color': DARK_GRAY},
])

# ============================================================
# 第4页：分析框架
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_title_bar(slide, '分析框架与方法论', 'ANALYTICAL FRAMEWORK')
add_page_number(slide, 4)

# 四个维度卡片
dims = [
    ('综合评价分析', '价格分布 · 好评率分布\n拥有者分布 · 发行趋势\n相关性矩阵', DARK_BLUE),
    ('游戏定价策略', '价格区间分布\n各区间好评率对比\n免费 vs 付费深度对比', MEDIUM_BLUE),
    ('游戏类型趋势', '类型市场份额\n类型好评率对比\n年度发行趋势', LIGHT_BLUE),
    ('玩家评价分析', '评价数 vs 好评率\nMetacritic 对比验证\n影响因素分组探究', ACCENT_ORANGE),
]

for i, (title, items, color) in enumerate(dims):
    left = Inches(0.6) + Inches(3.1) * i
    top = Inches(1.6)
    # 卡片背景
    add_rect(slide, left, top, Inches(2.8), Inches(3.5), BOX_BG)
    add_rect(slide, left, top, Inches(2.8), Inches(0.06), color)
    # 编号
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(1.0), top + Inches(0.4), Inches(0.8), Inches(0.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = str(i + 1)
    p.font.size = Pt(30)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    add_textbox(slide, left + Inches(0.15), top + Inches(1.5), Inches(2.5), Inches(0.5),
                title, font_size=18, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.15), top + Inches(2.1), Inches(2.5), Inches(1.2),
                items, font_size=12, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

# 方法标签
add_textbox(slide, Inches(0.8), Inches(5.4), Inches(11), Inches(0.4),
            '分析方法', font_size=16, color=DARK_BLUE, bold=True)
methods = ['描述性统计', 'Pearson相关性', '分组对比分析', '时间序列', '多维可视化']
for i, m in enumerate(methods):
    left = Inches(0.8) + Inches(2.4) * i
    shape = add_rect(slide, left, Inches(5.9), Inches(2.1), Inches(0.5), MEDIUM_BLUE)
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = m
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

add_textbox(slide, Inches(0.8), Inches(6.6), Inches(11.5), Inches(0.6),
            '▸ 多方法协同：描述统计提供概览 → 相关性发现关联 → 分组对比揭示结构 → 可视化验证假设 → 时间序列展示动态',
            font_size=12, color=MEDIUM_GRAY)

# ============================================================
# 第5页：方法选择优势（1）
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_title_bar(slide, '分析方法选择与优势', 'WHY THESE METHODS?')
add_page_number(slide, 5)

comparisons = [
    ('中位数 vs 均值', 'M', DARK_BLUE,
     ['游戏市场呈极端长尾分布，少数爆款严重拉高均值',
      '中位数对异常值鲁棒，反映"典型游戏"真实情况',
      '例：付费游戏中位价 $5.09，均价却被拉至 $9.83',
      '两者对比本身就传达了分布偏态程度的信息']),
    ('Pearson vs Spearman', 'P', MEDIUM_BLUE,
     ['大样本（数万条）下 Pearson 对非正态分布鲁棒',
      'Pearson 保留原始尺度信息，直观解释"量-效关系"',
      'Spearman 仅依赖排序，丢失数值大小的信息',
      '实际验证：两项系数在本数据集高度一致']),
    ('分组对比 vs 回归', 'G', LIGHT_BLUE,
     ['不受函数形式假设约束（无需预设线性/多项式）',
      '能捕捉非单调关系（如价格-好评率"倒U型"）',
      '分组统计量有直观业务含义，便于沟通',
      '回归分析可作为后续深化量化效应大小的方向']),
]

for i, (title, icon, color, points) in enumerate(comparisons):
    left = Inches(0.5) + Inches(4.2) * i
    # 卡片
    add_rect(slide, left, Inches(1.5), Inches(3.9), Inches(5.0), BOX_BG)
    add_rect(slide, left, Inches(1.5), Inches(3.9), Inches(0.06), color)
    # 图标
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.3), Inches(1.8), Inches(0.5), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = icon
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    add_textbox(slide, left + Inches(1.0), Inches(1.8), Inches(2.7), Inches(0.4),
                title, font_size=18, color=color, bold=True)

    lines = [{'text': '✓  ' + pt, 'size': 12, 'color': DARK_GRAY} for pt in points]
    add_multiline_textbox(slide, left + Inches(0.3), Inches(2.5), Inches(3.3), Inches(3.8), lines)

add_textbox(slide, Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.4),
            '核心原则：针对数据特征选择方法，而非套用"万能公式" — 偏态分布用中位数、非单调关系用分组对比、多维模式用可视化',
            font_size=12, color=ACCENT_ORANGE, bold=True)

# ============================================================
# 第6页：方法选择优势（2）
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_title_bar(slide, '可视化策略与多方法协同', 'VISUALIZATION & SYNERGY')
add_page_number(slide, 6)

# 左侧 — 可视化策略
add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.4),
            '可视化策略：多维优于单一', font_size=20, color=DARK_BLUE, bold=True)
add_multiline_textbox(slide, Inches(0.8), Inches(2.0), Inches(5.8), Inches(2.8), [
    {'text': '2×2 子图联动布局', 'size': 16, 'bold': True, 'color': MEDIUM_BLUE},
    {'text': '• 同一视野呈现四个互补视角', 'size': 13, 'color': DARK_GRAY},
    {'text': '• 子图空间邻近降低认知负荷', 'size': 13, 'color': DARK_GRAY},
    {'text': '• 直方图+散点图+箱线图+折线图组合', 'size': 13, 'color': DARK_GRAY},
    {'text': '', 'size': 8},
    {'text': '颜色编码第三变量', 'size': 16, 'bold': True, 'color': MEDIUM_BLUE},
    {'text': '• 二维平面上呈现三维信息', 'size': 13, 'color': DARK_GRAY},
    {'text': '• 例：价格-拥有者散点以颜色表好评率', 'size': 13, 'color': DARK_GRAY},
    {'text': '• 快速识别"高价格低销量高好评"群体', 'size': 13, 'color': DARK_GRAY},
    {'text': '', 'size': 8},
    {'text': '时间序列 vs 截面数据', 'size': 16, 'bold': True, 'color': MEDIUM_BLUE},
    {'text': '• 揭示增长拐点和增速差异', 'size': 13, 'color': DARK_GRAY},
    {'text': '• 发现市场饱和信号（2023年后放缓）', 'size': 13, 'color': DARK_GRAY},
    {'text': '• 识别外部事件冲击（2020疫情效应）', 'size': 13, 'color': DARK_GRAY},
])

# 右侧 — 递进验证体系
add_textbox(slide, Inches(7.3), Inches(1.5), Inches(5.5), Inches(0.4),
            '多方法协同：递进验证体系', font_size=20, color=DARK_BLUE, bold=True)

steps = [
    ('① 描述统计', '提供宏观概览', DARK_BLUE),
    ('② 相关分析', '发现潜在关联', MEDIUM_BLUE),
    ('③ 分组对比', '揭示结构差异', LIGHT_BLUE),
    ('④ 可视化', '呈现模式、验证假设', ACCENT_ORANGE),
    ('⑤ 时序分析', '展示动态演变', ACCENT_GREEN),
]
for i, (step, desc, color) in enumerate(steps):
    y = Inches(2.2) + Inches(0.85) * i
    # 箭头线（最后一项不加）
    if i < len(steps) - 1:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.2), y + Inches(0.55), Inches(0.04), Inches(0.3))
        shape.fill.solid()
        shape.fill.fore_color.rgb = MEDIUM_GRAY
        shape.line.fill.background()
    # 步骤卡片
    shape = add_rect(slide, Inches(7.5), y, Inches(5.2), Inches(0.55), color)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f'{step}  {desc}'
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

add_multiline_textbox(slide, Inches(7.3), Inches(6.4), Inches(5.5), Inches(0.8), [
    {'text': '▸ 优势：多角度交叉验证，降低"数据挖掘偏差"风险', 'size': 12, 'color': ACCENT_ORANGE, 'bold': True},
    {'text': '▸ 对比：单一方法（如仅回归）可能遗漏非单调关系或隐藏结构', 'size': 12, 'color': MEDIUM_GRAY},
])

# ============================================================
# 第7页：维度一 — 综合概览
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_title_bar(slide, '维度一：综合评价分析', 'DIMENSION 1: COMPREHENSIVE OVERVIEW')
add_page_number(slide, 7)

add_image_safe(slide, '01_综合概览.png', Inches(0.3), Inches(1.4), Inches(6.8))

# 右侧解读
add_textbox(slide, Inches(7.5), Inches(1.5), Inches(5.5), Inches(0.4),
            '核心发现', font_size=20, color=DARK_BLUE, bold=True)

findings = [
    {'text': '价格分布', 'size': 15, 'bold': True, 'color': MEDIUM_BLUE},
    {'text': '中位价 $5.09，高度右偏。60%+ 游戏定价$10以下，AAA级游戏极为稀缺', 'size': 12, 'color': DARK_GRAY},
    {'text': '', 'size': 6},
    {'text': '好评率分布', 'size': 15, 'bold': True, 'color': MEDIUM_BLUE},
    {'text': '中位好评率 81.3%，62.2% 游戏集中在70%-90%。极好评（95%+）极少', 'size': 12, 'color': DARK_GRAY},
    {'text': '', 'size': 6},
    {'text': '拥有者长尾', 'size': 15, 'bold': True, 'color': MEDIUM_BLUE},
    {'text': '中位拥有者仅 10,000 人。极少数爆款占据大量用户，市场高度集中', 'size': 12, 'color': DARK_GRAY},
    {'text': '', 'size': 6},
    {'text': '发行趋势', 'size': 15, 'bold': True, 'color': MEDIUM_BLUE},
    {'text': '2023年达峰值 15,542 款。2015年后爆发式增长，与 Steam Direct 降低门槛直接相关', 'size': 12, 'color': DARK_GRAY},
]
add_multiline_textbox(slide, Inches(7.5), Inches(2.0), Inches(5.5), Inches(5.0), findings)

# ============================================================
# 第8页：维度一 — 相关性矩阵
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_title_bar(slide, '维度一：关键指标相关性', 'CORRELATION MATRIX')
add_page_number(slide, 8)

add_image_safe(slide, '01_相关性热力图.png', Inches(0.3), Inches(1.4), Inches(6.5))

add_textbox(slide, Inches(7.3), Inches(1.5), Inches(5.5), Inches(0.4),
            '关键相关系数解读', font_size=20, color=DARK_BLUE, bold=True)

corr_findings = [
    {'text': '价格 vs 好评率: r = 0.046', 'size': 15, 'bold': True, 'color': ACCENT_RED},
    {'text': '几乎无关 — "贵的不一定好"，挑战"一分钱一分货"的直觉', 'size': 12, 'color': DARK_GRAY},
    {'text': '', 'size': 6},
    {'text': '拥有者 vs 峰值在线: r = 0.628', 'size': 15, 'bold': True, 'color': ACCENT_GREEN},
    {'text': '中等正相关 — 热门游戏拥有更高活跃度，但存在"买而不用"', 'size': 12, 'color': DARK_GRAY},
    {'text': '', 'size': 6},
    {'text': 'Metacritic vs Steam好评率: r = 0.475', 'size': 15, 'bold': True, 'color': MEDIUM_BLUE},
    {'text': '中等正相关 — 专业媒体与玩家审美大体一致但不完全重合', 'size': 12, 'color': DARK_GRAY},
    {'text': '', 'size': 6},
    {'text': '评价总数 vs 好评率: r = 0.019', 'size': 15, 'bold': True, 'color': ACCENT_ORANGE},
    {'text': '几乎无关 — 受欢迎不等于高评价，评价多寡不决定评分高低', 'size': 12, 'color': DARK_GRAY},
]
add_multiline_textbox(slide, Inches(7.3), Inches(2.0), Inches(5.5), Inches(4.5), corr_findings)

# ============================================================
# 第9页：维度二 — 定价策略
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_title_bar(slide, '维度二：游戏定价策略', 'DIMENSION 2: PRICING STRATEGY')
add_page_number(slide, 9)

add_image_safe(slide, '02_定价策略分析.png', Inches(0.3), Inches(1.4), Inches(6.5))

add_textbox(slide, Inches(7.3), Inches(1.5), Inches(5.5), Inches(0.4),
            '核心发现', font_size=20, color=DARK_BLUE, bold=True)

price_findings = [
    {'text': '"$5-$15"是"最佳甜蜜点"', 'size': 15, 'bold': True, 'color': ACCENT_GREEN},
    {'text': '该区间游戏中位好评率最高（83.2%），性价比预期管理最优', 'size': 12, 'color': DARK_GRAY},
    {'text': '', 'size': 6},
    {'text': '免费游戏：量大但口碑弱', 'size': 15, 'bold': True, 'color': ACCENT_ORANGE},
    {'text': '中位拥有者 75,000（付费的2倍+），但好评率仅 76.3%', 'size': 12, 'color': DARK_GRAY},
    {'text': '', 'size': 6},
    {'text': '付费游戏：精准但面窄', 'size': 15, 'bold': True, 'color': MEDIUM_BLUE},
    {'text': '好评率 82.1% 显著优于免费，但中位拥有者仅 35,000', 'size': 12, 'color': DARK_GRAY},
    {'text': '', 'size': 6},
    {'text': '高价游戏面临"期待落差"风险', 'size': 15, 'bold': True, 'color': ACCENT_RED},
    {'text': '高价区间 ($60+) 好评率仅 80.0%，玩家期望值高导致评价更严苛', 'size': 12, 'color': DARK_GRAY},
]
add_multiline_textbox(slide, Inches(7.3), Inches(2.0), Inches(5.5), Inches(4.8), price_findings)

# ============================================================
# 第10页：维度三 — 类型分析
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_title_bar(slide, '维度三：游戏类型趋势', 'DIMENSION 3: GENRE TRENDS')
add_page_number(slide, 10)

add_image_safe(slide, '03_类型分析.png', Inches(0.2), Inches(1.4), Inches(6.8))

add_textbox(slide, Inches(7.5), Inches(1.5), Inches(5.5), Inches(0.4),
            '核心发现', font_size=20, color=DARK_BLUE, bold=True)

genre_findings = [
    {'text': 'Indie 绝对主导', 'size': 15, 'bold': True, 'color': DARK_BLUE},
    {'text': '64,499 款独立游戏，占所有类型条目首位。Steam 是独立开发者首选平台', 'size': 12, 'color': DARK_GRAY},
    {'text': '', 'size': 6},
    {'text': 'Casual 好评率最高', 'size': 15, 'bold': True, 'color': ACCENT_GREEN},
    {'text': '休闲类以 83.8% 中位好评率领跑 Top 15。机制简单成熟、玩家预期合理', 'size': 12, 'color': DARK_GRAY},
    {'text': '', 'size': 6},
    {'text': 'Free to Play 评分最低', 'size': 15, 'bold': True, 'color': ACCENT_RED},
    {'text': '仅 75.0% — 免费玩家基数大但投入度低、容忍度低', 'size': 12, 'color': DARK_GRAY},
    {'text': '', 'size': 6},
    {'text': '类型格局：三巨头 + 长尾', 'size': 15, 'bold': True, 'color': MEDIUM_BLUE},
    {'text': 'Indie + Action + Adventure 覆盖 14万+ 条目，其余类型构成丰富长尾', 'size': 12, 'color': DARK_GRAY},
    {'text': '', 'size': 6},
    {'text': 'Design & Illustration 定价最高 ($25.72)，工具类软件价格普遍偏高', 'size': 12, 'color': MEDIUM_GRAY},
]
add_multiline_textbox(slide, Inches(7.5), Inches(2.0), Inches(5.5), Inches(5.0), genre_findings)

# ============================================================
# 第11页：维度四 — 评价分析
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_title_bar(slide, '维度四：玩家评价分析', 'DIMENSION 4: PLAYER REVIEWS')
add_page_number(slide, 11)

add_image_safe(slide, '04_评价分析.png', Inches(0.2), Inches(1.4), Inches(6.8))

add_textbox(slide, Inches(7.5), Inches(1.5), Inches(5.5), Inches(0.4),
            '核心发现', font_size=20, color=DARK_BLUE, bold=True)

review_findings = [
    {'text': '评价数量 ≠ 评价质量', 'size': 15, 'bold': True, 'color': ACCENT_ORANGE},
    {'text': 'r=0.019，两者几乎无关。高评价量游戏的好评率方差极大（20%-100%）', 'size': 12, 'color': DARK_GRAY},
    {'text': '', 'size': 6},
    {'text': 'Metacritic 与 Steam 有中等一致性', 'size': 15, 'bold': True, 'color': MEDIUM_BLUE},
    {'text': 'r=0.475 — 专业媒体与玩家存在审美一致，但远非完全重合', 'size': 12, 'color': DARK_GRAY},
    {'text': '', 'size': 6},
    {'text': '多平台支持提升好评率', 'size': 15, 'bold': True, 'color': ACCENT_GREEN},
    {'text': '支持3平台 > 2平台 > 1平台（83.2% > 82.1% > 80.5%）', 'size': 12, 'color': DARK_GRAY},
    {'text': '', 'size': 6},
    {'text': 'DLC 与成就系统的正向效应', 'size': 15, 'bold': True, 'color': ACCENT_GREEN},
    {'text': '有DLC游戏好评率高2.3%，丰富成就系统 (+100) 好评率高3.3%', 'size': 12, 'color': DARK_GRAY},
]
add_multiline_textbox(slide, Inches(7.5), Inches(2.0), Inches(5.5), Inches(4.5), review_findings)

# ============================================================
# 第12页：评价影响因素对比
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_title_bar(slide, '影响好评率的关键因素', 'FACTORS AFFECTING REVIEW SCORES')
add_page_number(slide, 12)

# 三个对比组
factors = [
    ('平台支持数量', [
        ('1个平台', '80.5%', MEDIUM_BLUE),
        ('2个平台', '82.1%', LIGHT_BLUE),
        ('3个平台', '83.2%', ACCENT_GREEN),
    ]),
    ('DLC 情况', [
        ('无 DLC', '80.5%', MEDIUM_BLUE),
        ('有 DLC', '82.8%', ACCENT_GREEN),
    ]),
    ('成就系统丰富度', [
        ('≤100 成就', '80.2%', MEDIUM_BLUE),
        ('>100 成就', '83.5%', ACCENT_GREEN),
    ]),
]

for i, (factor_name, items) in enumerate(factors):
    left = Inches(0.5) + Inches(4.2) * i
    add_textbox(slide, left, Inches(1.5), Inches(3.8), Inches(0.5),
                factor_name, font_size=18, color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

    for j, (label, value, color) in enumerate(items):
        y = Inches(2.2) + Inches(1.3) * j
        add_rect(slide, left + Inches(0.15), y, Inches(3.5), Inches(0.95), BOX_BG)
        add_textbox(slide, left + Inches(0.3), y + Inches(0.1), Inches(2.2), Inches(0.4),
                    label, font_size=16, color=DARK_GRAY, bold=True)
        add_textbox(slide, left + Inches(0.3), y + Inches(0.5), Inches(2.2), Inches(0.4),
                    value, font_size=24, color=color, bold=True)

    # 趋势箭头
    if len(items) >= 2 and float(items[-1][1].replace('%', '')) > float(items[0][1].replace('%', '')):
        diff = float(items[-1][1].replace('%', '')) - float(items[0][1].replace('%', ''))
        add_textbox(slide, left + Inches(0.15), Inches(2.2) + Inches(1.3) * len(items) + Inches(0.1),
                    Inches(3.5), Inches(0.4),
                    f'↑ 提升 {diff:.1f} 个百分点', font_size=13, color=ACCENT_GREEN,
                    bold=True, alignment=PP_ALIGN.CENTER)

# 底部解读
add_textbox(slide, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.9),
            '解读：持续运营投入（多平台移植、DLC 开发、成就系统设计）能够显著提升玩家满意度。'
            '这不仅是品质信号，也筛选出更忠实的用户群体，形成良性循环。',
            font_size=14, color=ACCENT_ORANGE)

# ============================================================
# 第13页：高评价 vs 低评价
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_title_bar(slide, '高评价 vs 低评价游戏对比', 'HIGH-RATED vs LOW-RATED GAMES')
add_page_number(slide, 13)

# 对比表
add_textbox(slide, Inches(0.8), Inches(1.5), Inches(3), Inches(0.4),
            '高评价游戏 (≥90%)', font_size=20, color=ACCENT_GREEN, bold=True)
add_textbox(slide, Inches(7.3), Inches(1.5), Inches(3), Inches(0.4),
            '低评价游戏 (<60%)', font_size=20, color=ACCENT_RED, bold=True)

# 左侧 - 高评价
add_rect(slide, Inches(0.8), Inches(2.1), Inches(5.5), Inches(3.5), BOX_BG)
high_items = [
    {'text': '6,246 款', 'size': 36, 'bold': True, 'color': ACCENT_GREEN},
    {'text': '占评价样本的 26.1%', 'size': 14, 'color': DARK_GRAY},
    {'text': '', 'size': 10},
    {'text': '✓ 平均价格较低', 'size': 14, 'color': DARK_GRAY},
    {'text': '✓ 中位拥有者较广泛', 'size': 14, 'color': DARK_GRAY},
    {'text': '✓ 常见类型：Casual、Indie', 'size': 14, 'color': DARK_GRAY},
    {'text': '✓ 品质稳定，预期管理良好', 'size': 14, 'color': DARK_GRAY},
]
add_multiline_textbox(slide, Inches(1.2), Inches(2.3), Inches(5), Inches(3), high_items)

# 右侧 - 低评价
add_rect(slide, Inches(7.3), Inches(2.1), Inches(5.5), Inches(3.5), BOX_BG)
low_items = [
    {'text': '3,518 款', 'size': 36, 'bold': True, 'color': ACCENT_RED},
    {'text': '占评价样本的 14.7%', 'size': 14, 'color': DARK_GRAY},
    {'text': '', 'size': 10},
    {'text': '✗ 平均价格较高', 'size': 14, 'color': DARK_GRAY},
    {'text': '✗ 中位拥有者较少', 'size': 14, 'color': DARK_GRAY},
    {'text': '✗ Free to Play、Early Access 比例偏高', 'size': 14, 'color': DARK_GRAY},
    {'text': '✗ 期待落差大，预期管理不足', 'size': 14, 'color': DARK_GRAY},
]
add_multiline_textbox(slide, Inches(7.7), Inches(2.3), Inches(5), Inches(3), low_items)

# 底部建议
add_textbox(slide, Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.8),
            '⚠ 警示：Early Access（抢先体验）需谨慎 — 虽然能早期获取资金和反馈，但品质不达标时容易沦为低评价重灾区',
            font_size=14, color=ACCENT_ORANGE, bold=True)

# ============================================================
# 第14页：核心发现总结
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_title_bar(slide, '核心发现总结', 'KEY FINDINGS')
add_page_number(slide, 14)

key_findings = [
    ('01', '市场高度长尾化', '中位拥有者仅 10,000 人，极少数爆款占据大量用户。独立开发者需在成本控制和利基市场中找到平衡', DARK_BLUE),
    ('02', '定价与质量脱钩', 'r=0.046，价格几乎不影响评价。Steam 平台为优质低价游戏提供了公平竞争的环境', MEDIUM_BLUE),
    ('03', '付费体验优于免费', '付费好评率 82.1% vs 免费 76.3%。但免费游戏用户覆盖面是付费的 2 倍以上', LIGHT_BLUE),
    ('04', '$5-$15 是甜蜜点', '该区间好评率最高（83.2%），性价比预期管理最优。独立游戏定价的黄金区间', ACCENT_GREEN),
    ('05', 'Indie 为王、Casual 为优', '独立游戏数量压倒性领先。休闲类好评率最高（83.8%），用户体验最佳', ACCENT_ORANGE),
    ('06', '持续运营带来正反馈', '多平台 + DLC + 成就系统 = 更高好评率。玩家认可开发者的长期投入', DARK_BLUE),
]

for i, (num, title, desc, color) in enumerate(key_findings):
    row = i // 2
    col = i % 2
    left = Inches(0.5) + Inches(6.3) * col
    top = Inches(1.5) + Inches(1.8) * row

    # 编号圆
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.1), top + Inches(0.15), Inches(0.5), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    add_textbox(slide, left + Inches(0.8), top, Inches(5), Inches(0.4),
                title, font_size=18, color=color, bold=True)
    add_textbox(slide, left + Inches(0.8), top + Inches(0.5), Inches(5.2), Inches(0.9),
                desc, font_size=12, color=DARK_GRAY)

# ============================================================
# 第15页：方法论回顾
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_title_bar(slide, '方法论特色总结', 'METHODOLOGY HIGHLIGHTS')
add_page_number(slide, 15)

methods_highlights = [
    ('数据思维', [
        '从"提问题"出发，而非从数据出发',
        '四个分析维度基于真实市场关切设计',
        '每个发现都能回答一个实际决策问题',
    ]),
    ('方法适配', [
        '偏态分布 → 中位数替代均值',
        '非单调关系 → 分组对比替代单一回归',
        '多维信息 → 颜色编码+多子图联动',
    ]),
    ('递进验证', [
        '描述统计 → 相关分析 → 分组对比 → 可视化',
        '多方法交叉印证，降低伪发现风险',
        '单一方法无法捕捉的"倒U型"关系被成功揭示',
    ]),
    ('可操作性', [
        '每个分析维度输出可执行的建议',
        '定价"甜蜜点"可直接指导商业决策',
        '开发者可基于结论优化产品策略',
    ]),
]

for i, (title, points) in enumerate(methods_highlights):
    left = Inches(0.4) + Inches(3.2) * i
    add_rect(slide, left, Inches(1.6), Inches(2.9), Inches(4.5), BOX_BG)
    add_rect(slide, left, Inches(1.6), Inches(2.9), Inches(0.06), [DARK_BLUE, MEDIUM_BLUE, LIGHT_BLUE, ACCENT_ORANGE][i])

    add_textbox(slide, left + Inches(0.15), Inches(1.9), Inches(2.6), Inches(0.4),
                title, font_size=18, color=[DARK_BLUE, MEDIUM_BLUE, LIGHT_BLUE, ACCENT_ORANGE][i],
                bold=True, alignment=PP_ALIGN.CENTER)

    pts = [{'text': '• ' + p, 'size': 12, 'color': DARK_GRAY} for p in points]
    add_multiline_textbox(slide, left + Inches(0.2), Inches(2.5), Inches(2.5), Inches(3.5), pts)

# ============================================================
# 第16页：建议与展望
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_title_bar(slide, '对开发者的建议与未来展望', 'RECOMMENDATIONS & OUTLOOK')
add_page_number(slide, 16)

add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.4),
            '对 Steam 开发者的五项建议', font_size=20, color=DARK_BLUE, bold=True)

recs = [
    ('💲', '定价策略', '优先考虑 $5-$15 中低价位，兼顾覆盖和口碑', MEDIUM_BLUE),
    ('🎯', '商业选择', '用户增长选免费，品牌口碑选付费', LIGHT_BLUE),
    ('🔄', '持续运营', '开发 DLC、成就系统，支持多平台', ACCENT_GREEN),
    ('🎮', '类型策略', 'Indie/Casual 赛道友好，但需差异化突围', ACCENT_ORANGE),
    ('⚠️', 'Early Access', '核心玩法完善后再用，避免低评价陷阱', ACCENT_RED),
]

for i, (icon, title, desc, color) in enumerate(recs):
    y = Inches(2.2) + Inches(0.55) * i
    add_icon_bullet(slide, Inches(0.8), y, Inches(5.5), Inches(0.5),
                    icon, title, desc, color)

# 右侧 — 展望
add_textbox(slide, Inches(7.3), Inches(1.5), Inches(5.5), Inches(0.4),
            '研究不足与未来方向', font_size=20, color=DARK_BLUE, bold=True)

outlook = [
    {'text': '数据层面', 'size': 16, 'bold': True, 'color': MEDIUM_BLUE},
    {'text': '• 当前为静态快照，缺乏时间序列变化', 'size': 13, 'color': DARK_GRAY},
    {'text': '• 拥有者数据为区间估算，精度有限', 'size': 13, 'color': DARK_GRAY},
    {'text': '', 'size': 8},
    {'text': '方法层面', 'size': 16, 'bold': True, 'color': MEDIUM_BLUE},
    {'text': '• 可引入多元回归量化效应大小', 'size': 13, 'color': DARK_GRAY},
    {'text': '• NLP 分析评论文本情感更为深入', 'size': 13, 'color': DARK_GRAY},
    {'text': '• 因果推断探索"价格→评价"的因果路径', 'size': 13, 'color': DARK_GRAY},
    {'text': '', 'size': 8},
    {'text': '扩展方向', 'size': 16, 'bold': True, 'color': MEDIUM_BLUE},
    {'text': '• 地域/文化差异对游戏偏好的影响', 'size': 13, 'color': DARK_GRAY},
    {'text': '• 标签组合分析（Tags 共现网络）', 'size': 13, 'color': DARK_GRAY},
    {'text': '• 预测模型：什么特征造就"爆款"？', 'size': 13, 'color': DARK_GRAY},
]
add_multiline_textbox(slide, Inches(7.3), Inches(2.0), Inches(5.5), Inches(5.0), outlook)

# ============================================================
# 第17页：致谢
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)

add_rect(slide, Inches(5.6), Inches(3.0), Inches(2.0), Inches(0.06), ACCENT_ORANGE)

add_textbox(slide, Inches(1.5), Inches(3.3), Inches(10), Inches(1.0),
            '感谢聆听', font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1.5), Inches(4.5), Inches(10), Inches(0.6),
            'Q & A', font_size=28, color=RGBColor(0xBD, 0xC3, 0xC7), alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1.5), Inches(5.8), Inches(10), Inches(0.4),
            '数据思维与人工智能 · 课程大作业', font_size=14, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# 保存
# ============================================================
output_path = 'Steam游戏市场数据分析_演讲.pptx'
prs.save(output_path)
print(f'PPT 已生成: {output_path}')
print(f'共 {len(prs.slides)} 页幻灯片')
